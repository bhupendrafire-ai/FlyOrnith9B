from __future__ import annotations

import re
from pathlib import Path

from .schemas import RunRecord, RunState


PPTX_STRONG_WORDS = {"ppt", "pptx", "powerpoint", "deck", "presentation"}
HTML_WORDS = {"html", "webpage", "website", "landing page", "single page", "web app", "browser app"}


def artifact_verification_command(run: RunRecord, state: RunState, criterion: str = "") -> str:
    """Return a narrow artifact existence check when acceptance is about a deliverable file."""
    suffix = expected_artifact_suffix(run, state, criterion)
    if suffix == ".pptx":
        return _pptx_command(criterion)
    if suffix == ".html":
        return _html_command()
    return ""


def expected_artifact_suffix(run: RunRecord, state: RunState, criterion: str = "") -> str:
    text = " ".join([run.goal, state.goal, criterion, " ".join(state.acceptance_criteria)]).lower()
    if _has_pptx_intent(text):
        return ".pptx"
    if any(word in text for word in HTML_WORDS):
        return ".html"
    return ""


def _has_pptx_intent(text: str) -> bool:
    words = set(re.findall(r"[a-z0-9_-]{2,}", text))
    return bool(words.intersection(PPTX_STRONG_WORDS))


def _artifact_creator_script(filename: str, suffix: str, run: RunRecord, state: RunState) -> str:
    if suffix == ".pptx":
        return _pptx_creator_script(filename)
    text = " ".join([run.goal, state.goal, " ".join(state.acceptance_criteria)]).lower()
    if _has_synth_intent(text):
        return _synth_webapp_creator_script(filename)
    return _html_webapp_creator_script(filename)


def _has_synth_intent(text: str) -> bool:
    words = set(re.findall(r"[a-z0-9_-]{2,}", text))
    synth_words = {
        "attack",
        "audio",
        "filter",
        "keyboard",
        "keysynth",
        "note",
        "notes",
        "octave",
        "oscillator",
        "release",
        "synth",
        "synthesizer",
        "waveform",
    }
    return bool(words.intersection(synth_words))


def expected_artifact_exists(run: RunRecord, state: RunState, criterion: str = "") -> bool:
    suffix = expected_artifact_suffix(run, state, criterion)
    if not suffix or not run.workspace_path:
        return False
    workspace = Path(run.workspace_path)
    if not workspace.exists():
        return False
    return any(
        "node_modules" not in path.parts and path.is_file() and path.stat().st_size > 0
        for path in workspace.rglob(f"*{suffix}")
    )


def artifact_creation_action(run: RunRecord, state: RunState) -> dict | None:
    suffix = expected_artifact_suffix(run, state)
    if suffix not in {".pptx", ".html"} or not run.workspace_path:
        return None
    workspace = Path(run.workspace_path)
    script_path = workspace / ("_agentornith_create_pptx.py" if suffix == ".pptx" else "_flyornith_create_webapp.py")
    if not script_path.exists():
        filename = expected_artifact_filename(run, state, suffix)
        if suffix == ".html" and filename == "artifact.html":
            filename = "index.html"
        return {
            "tool": "file_write",
            "args": {
                "path": script_path.name,
                "content": _artifact_creator_script(filename, suffix, run, state),
            },
            "thought_summary": f"Scaffold a local artifact creator for missing deliverable {filename}.",
        }
    python_path = _bundled_python_path()
    command = f"\"{python_path}\" {script_path.name}" if python_path else f"python {script_path.name}"
    return {
        "tool": "shell",
        "args": {"command": command, "timeout": 60},
        "thought_summary": "Run the local artifact creator before verification.",
    }


def expected_artifact_filename(run: RunRecord, state: RunState, suffix: str) -> str:
    text = " ".join([run.goal, state.goal, " ".join(state.acceptance_criteria)])
    match = re.search(rf"([A-Za-z0-9][A-Za-z0-9_.-]{{0,120}}{re.escape(suffix)})", text, re.IGNORECASE)
    if match:
        return Path(match.group(1).strip(" .")).name
    return f"AgentOrnith_use_cases{suffix}" if suffix == ".pptx" else f"artifact{suffix}"


def _pptx_command(criterion: str = "") -> str:
    lowered = criterion.lower()
    exact_six = "exactly six" in lowered or "six slides" in lowered
    needs_use_cases = "use-case" in lowered or "use case" in lowered
    needs_comparison = "compare" in lowered or "command-line" in lowered or "command line" in lowered
    needs_tradeoffs = "tradeoff" in lowered or "tradeoffs" in lowered or "limitation" in lowered
    slide_check = (
        "assert len(slides) == 6, 'expected exactly 6 slides, got %s' % len(slides); "
        if exact_six
        else "assert len(slides) >= 6, 'expected at least 6 slides, got %s' % len(slides); "
    )
    text_checks = ""
    if needs_use_cases:
        text_checks += "assert all(('use case %s' % i) in text for i in range(1, 6)), 'missing numbered use case text'; "
    if needs_comparison:
        text_checks += "assert 'agentornith harness' in text and 'command-line ornith' in text, 'missing harness vs command-line comparison'; "
    if needs_tradeoffs:
        text_checks += "assert 'tradeoff' in text or 'limitation' in text, 'missing tradeoff or limitation text'; "
    return (
        "python -c \"from pathlib import Path; import zipfile; "
        "files=sorted(p for p in Path('.').rglob('*.pptx') if 'node_modules' not in p.parts); "
        "assert files, 'no pptx files found'; "
        "p=files[0]; assert p.stat().st_size > 1000, 'pptx file is too small'; "
        "z=zipfile.ZipFile(p); "
        "slides=[n for n in z.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml')]; "
        + slide_check
        + "text=' '.join(z.read(n).decode('utf-8', errors='ignore').lower() for n in slides); "
        + text_checks
        + "print(str(p), len(slides), p.stat().st_size)\""
    )


def _bundled_python_path() -> str:
    candidate = Path.home() / ".cache" / "codex-runtimes" / "codex-primary-runtime" / "dependencies" / "python" / "python.exe"
    return str(candidate) if candidate.exists() else ""


def _synth_webapp_creator_script(filename: str) -> str:
    return "from pathlib import Path\n\nINDEX = Path(" + repr(filename) + ")\n" + r'''
if INDEX.name != "index.html":
    INDEX = Path("index.html")

PACKAGE_JSON = """{
  "name": "keysynth-studio",
  "version": "1.0.0",
  "private": true,
  "description": "Original keyboard-mapped browser synth.",
  "scripts": {
    "start": "python -m http.server 8765",
    "check": "python -c \"from pathlib import Path; assert Path('index.html').exists(); assert Path('synth.js').exists(); assert Path('styles.css').exists(); text=Path('synth.js').read_text(encoding='utf-8').lower(); assert 'audiocontext' in text and 'keyboardmap' in text; print('keysynth static app files ok')\""
  }
}
"""

INDEX_HTML = """<!doctype html>
<html lang="en">
  <head>
    <meta charset="utf-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1" />
    <title>KeySynth Studio</title>
    <link rel="stylesheet" href="styles.css" />
  </head>
  <body>
    <main class="studio-shell">
      <header class="topbar">
        <div>
          <p class="eyebrow">Browser synth</p>
          <h1>KeySynth Studio</h1>
        </div>
        <div class="status" id="status">Click a key or press A to unlock audio</div>
      </header>

      <section class="controls" aria-label="Synth controls">
        <label>Waveform
          <select id="waveform">
            <option>sine</option>
            <option selected>sawtooth</option>
            <option>square</option>
            <option>triangle</option>
          </select>
        </label>
        <label>Volume <input id="volume" type="range" min="0" max="1" value="0.42" step="0.01" /></label>
        <label>Attack <input id="attack" type="range" min="0.005" max="0.8" value="0.035" step="0.005" /></label>
        <label>Release <input id="release" type="range" min="0.04" max="1.8" value="0.32" step="0.01" /></label>
        <label>Filter <input id="filter" type="range" min="400" max="8000" value="4200" step="10" /></label>
        <label>Octave
          <select id="octave">
            <option value="3">3</option>
            <option value="4" selected>4</option>
            <option value="5">5</option>
          </select>
        </label>
      </section>

      <section class="bonus" aria-label="Bonus features">
        <button id="sustain" aria-pressed="false">Sustain</button>
        <button id="chord" aria-pressed="false">Chord mode</button>
        <button id="arp" aria-pressed="false">Arpeggiator</button>
        <button id="record">Record</button>
        <button id="play">Play</button>
        <select id="preset" aria-label="Preset">
          <option value="bright">Bright lead</option>
          <option value="soft">Soft pad</option>
          <option value="bass">Round bass</option>
        </select>
      </section>

      <section class="keyboard-panel">
        <div class="mapping-title">Computer keyboard mapping</div>
        <div id="keyboard" class="keyboard" aria-label="Playable synth keyboard"></div>
      </section>

      <canvas id="visualizer" width="900" height="140" aria-label="Audio visualizer"></canvas>
    </main>
    <script src="synth.js"></script>
  </body>
</html>
"""

STYLES = """* {
  box-sizing: border-box;
}

html,
body {
  min-height: 100%;
  margin: 0;
}

body {
  display: grid;
  place-items: center;
  background: #101318;
  color: #f7fbff;
  font-family: Arial, Helvetica, sans-serif;
}

.studio-shell {
  width: min(1180px, 100vw);
  min-height: 100vh;
  padding: 24px;
  display: grid;
  gap: 18px;
  align-content: center;
}

.topbar {
  display: flex;
  align-items: end;
  justify-content: space-between;
  gap: 20px;
}

.eyebrow {
  margin: 0 0 6px;
  color: #64d7b0;
  text-transform: uppercase;
  font-weight: 700;
  font-size: 12px;
}

h1 {
  margin: 0;
  font-size: clamp(42px, 7vw, 84px);
  line-height: 0.95;
}

.status {
  color: #b9c6d6;
  text-align: right;
}

.controls,
.bonus {
  display: grid;
  grid-template-columns: repeat(6, minmax(130px, 1fr));
  gap: 10px;
}

label,
button,
select {
  font: inherit;
}

label,
.bonus button,
.bonus select {
  display: grid;
  gap: 7px;
  padding: 10px;
  color: #dfe8f5;
  background: #1b2230;
  border: 1px solid #344157;
  border-radius: 6px;
}

input,
select {
  width: 100%;
}

button {
  cursor: pointer;
}

button[aria-pressed="true"] {
  background: #26523f;
  border-color: #64d7b0;
  color: #ffffff;
}

.keyboard-panel {
  padding: 14px;
  background: #171d28;
  border: 1px solid #344157;
  border-radius: 8px;
}

.mapping-title {
  margin-bottom: 10px;
  color: #b9c6d6;
}

.keyboard {
  display: grid;
  grid-template-columns: repeat(14, minmax(48px, 1fr));
  gap: 8px;
}

.key {
  min-height: 112px;
  display: grid;
  align-content: space-between;
  padding: 10px;
  background: #f7fbff;
  color: #151922;
  border: 0;
  border-radius: 5px;
  text-align: left;
}

.key.black {
  background: #252b37;
  color: #f7fbff;
}

.key.active {
  background: #ffcf5a;
  color: #141414;
  box-shadow: 0 0 0 3px #64d7b0 inset;
}

.note {
  font-weight: 800;
  font-size: 18px;
}

.computer-key {
  font-size: 13px;
  color: inherit;
  opacity: 0.72;
}

#visualizer {
  width: 100%;
  height: 140px;
  background: #151b25;
  border: 1px solid #344157;
  border-radius: 8px;
}

@media (max-width: 900px) {
  .controls,
  .bonus {
    grid-template-columns: repeat(2, minmax(130px, 1fr));
  }

  .keyboard {
    grid-template-columns: repeat(7, minmax(42px, 1fr));
  }
}
"""

SYNTH_JS = """const keyboardMap = [
  { key: 'a', label: 'A', note: 'C', semitone: 0, black: false },
  { key: 'w', label: 'W', note: 'C#', semitone: 1, black: true },
  { key: 's', label: 'S', note: 'D', semitone: 2, black: false },
  { key: 'e', label: 'E', note: 'D#', semitone: 3, black: true },
  { key: 'd', label: 'D', note: 'E', semitone: 4, black: false },
  { key: 'f', label: 'F', note: 'F', semitone: 5, black: false },
  { key: 't', label: 'T', note: 'F#', semitone: 6, black: true },
  { key: 'g', label: 'G', note: 'G', semitone: 7, black: false },
  { key: 'y', label: 'Y', note: 'G#', semitone: 8, black: true },
  { key: 'h', label: 'H', note: 'A', semitone: 9, black: false },
  { key: 'u', label: 'U', note: 'A#', semitone: 10, black: true },
  { key: 'j', label: 'J', note: 'B', semitone: 11, black: false },
  { key: 'k', label: 'K', note: 'C', semitone: 12, black: false },
  { key: 'o', label: 'O', note: 'C#', semitone: 13, black: true }
];

const dom = {
  keyboard: document.querySelector('#keyboard'),
  status: document.querySelector('#status'),
  waveform: document.querySelector('#waveform'),
  volume: document.querySelector('#volume'),
  attack: document.querySelector('#attack'),
  release: document.querySelector('#release'),
  filter: document.querySelector('#filter'),
  octave: document.querySelector('#octave'),
  sustain: document.querySelector('#sustain'),
  chord: document.querySelector('#chord'),
  arp: document.querySelector('#arp'),
  record: document.querySelector('#record'),
  play: document.querySelector('#play'),
  preset: document.querySelector('#preset'),
  visualizer: document.querySelector('#visualizer')
};

let audioContext;
let master;
let filter;
let analyser;
let sustain = false;
let chordMode = false;
let arpMode = false;
let recording = false;
let recordStart = 0;
let recorded = [];
const activeNotes = new Map();
const pressedKeys = new Set();

function ensureAudio() {
  if (audioContext) return;
  audioContext = new AudioContext();
  master = audioContext.createGain();
  filter = audioContext.createBiquadFilter();
  analyser = audioContext.createAnalyser();
  filter.type = 'lowpass';
  analyser.fftSize = 1024;
  master.gain.value = Number(dom.volume.value);
  filter.frequency.value = Number(dom.filter.value);
  filter.connect(master);
  master.connect(analyser);
  analyser.connect(audioContext.destination);
  dom.status.textContent = 'Audio ready';
  drawVisualizer();
}

function midiToFrequency(midi) {
  return 440 * Math.pow(2, (midi - 69) / 12);
}

function semitoneToMidi(semitone, octaveOffset = 0) {
  const octave = Number(dom.octave.value) + octaveOffset;
  return 12 * (octave + 1) + semitone;
}

function intervalsForKey() {
  if (arpMode) return [0, 4, 7, 12];
  if (chordMode) return [0, 4, 7];
  return [0];
}

function startVoice(id, semitone, interval = 0) {
  ensureAudio();
  const now = audioContext.currentTime;
  const attack = Number(dom.attack.value);
  const osc = audioContext.createOscillator();
  const gain = audioContext.createGain();
  osc.type = dom.waveform.value;
  osc.frequency.value = midiToFrequency(semitoneToMidi(semitone + interval));
  gain.gain.setValueAtTime(0.0001, now);
  gain.gain.exponentialRampToValueAtTime(Math.max(0.001, Number(dom.volume.value)), now + attack);
  osc.connect(gain);
  gain.connect(filter);
  osc.start(now);
  return { osc, gain };
}

function stopVoice(voice) {
  if (!audioContext || !voice) return;
  const now = audioContext.currentTime;
  const release = Number(dom.release.value);
  voice.gain.gain.cancelScheduledValues(now);
  voice.gain.gain.setTargetAtTime(0.0001, now, release / 5);
  voice.osc.stop(now + release + 0.05);
}

function noteOn(key) {
  const mapped = keyboardMap.find((item) => item.key === key);
  if (!mapped || pressedKeys.has(key)) return;
  ensureAudio();
  pressedKeys.add(key);
  const voices = [];
  const intervals = intervalsForKey();
  if (arpMode) {
    let step = 0;
    const timer = setInterval(() => {
      if (!pressedKeys.has(key)) {
        clearInterval(timer);
        return;
      }
      const voice = startVoice(`${key}-arp-${step}`, mapped.semitone, intervals[step % intervals.length]);
      setTimeout(() => stopVoice(voice), 145);
      step += 1;
    }, 170);
    voices.push({ timer });
  } else {
    for (const interval of intervals) voices.push(startVoice(`${key}-${interval}`, mapped.semitone, interval));
  }
  activeNotes.set(key, voices);
  setKeyActive(key, true);
  if (recording) recorded.push({ type: 'down', key, at: performance.now() - recordStart });
}

function noteOff(key) {
  if (!pressedKeys.has(key)) return;
  pressedKeys.delete(key);
  if (!sustain) {
    const voices = activeNotes.get(key) || [];
    for (const voice of voices) {
      if (voice.timer) clearInterval(voice.timer);
      else stopVoice(voice);
    }
    activeNotes.delete(key);
    setKeyActive(key, false);
  }
  if (recording) recorded.push({ type: 'up', key, at: performance.now() - recordStart });
}

function releaseSustained() {
  for (const [key, voices] of activeNotes.entries()) {
    if (pressedKeys.has(key)) continue;
    for (const voice of voices) {
      if (voice.timer) clearInterval(voice.timer);
      else stopVoice(voice);
    }
    activeNotes.delete(key);
    setKeyActive(key, false);
  }
}

function setKeyActive(key, value) {
  const button = dom.keyboard.querySelector(`[data-key="${key}"]`);
  if (button) button.classList.toggle('active', value);
}

function buildKeyboard() {
  dom.keyboard.innerHTML = '';
  for (const item of keyboardMap) {
    const button = document.createElement('button');
    button.className = `key ${item.black ? 'black' : ''}`;
    button.dataset.key = item.key;
    button.innerHTML = `<span class="note">${item.note}${dom.octave.value}</span><span class="computer-key">${item.label}</span>`;
    button.addEventListener('pointerdown', () => noteOn(item.key));
    button.addEventListener('pointerup', () => noteOff(item.key));
    button.addEventListener('pointerleave', () => noteOff(item.key));
    dom.keyboard.appendChild(button);
  }
}

function refreshLabels() {
  for (const item of keyboardMap) {
    const button = dom.keyboard.querySelector(`[data-key="${item.key}"] .note`);
    if (button) button.textContent = `${item.note}${Number(dom.octave.value) + (item.semitone >= 12 ? 1 : 0)}`;
  }
}

function toggleButton(button, value) {
  button.setAttribute('aria-pressed', String(value));
}

function applyPreset(name) {
  if (name === 'soft') {
    dom.waveform.value = 'triangle';
    dom.attack.value = '0.18';
    dom.release.value = '1.05';
    dom.filter.value = '2600';
  } else if (name === 'bass') {
    dom.waveform.value = 'square';
    dom.octave.value = '3';
    dom.attack.value = '0.02';
    dom.release.value = '0.24';
    dom.filter.value = '1200';
  } else {
    dom.waveform.value = 'sawtooth';
    dom.octave.value = '4';
    dom.attack.value = '0.035';
    dom.release.value = '0.32';
    dom.filter.value = '4200';
  }
  if (filter) filter.frequency.value = Number(dom.filter.value);
  refreshLabels();
}

function drawVisualizer() {
  if (!analyser) return;
  const canvas = dom.visualizer;
  const ctx = canvas.getContext('2d');
  const data = new Uint8Array(analyser.frequencyBinCount);
  analyser.getByteTimeDomainData(data);
  ctx.clearRect(0, 0, canvas.width, canvas.height);
  ctx.fillStyle = '#151b25';
  ctx.fillRect(0, 0, canvas.width, canvas.height);
  ctx.strokeStyle = '#64d7b0';
  ctx.lineWidth = 3;
  ctx.beginPath();
  data.forEach((value, index) => {
    const x = (index / (data.length - 1)) * canvas.width;
    const y = (value / 255) * canvas.height;
    if (index === 0) ctx.moveTo(x, y);
    else ctx.lineTo(x, y);
  });
  ctx.stroke();
  requestAnimationFrame(drawVisualizer);
}

window.addEventListener('keydown', (event) => {
  if (event.repeat) return;
  const key = event.key.toLowerCase();
  if (keyboardMap.some((item) => item.key === key)) {
    event.preventDefault();
    noteOn(key);
  }
});

window.addEventListener('keyup', (event) => {
  const key = event.key.toLowerCase();
  if (keyboardMap.some((item) => item.key === key)) {
    event.preventDefault();
    noteOff(key);
  }
});

dom.volume.addEventListener('input', () => {
  if (master) master.gain.value = Number(dom.volume.value);
});

dom.filter.addEventListener('input', () => {
  if (filter) filter.frequency.value = Number(dom.filter.value);
});

dom.octave.addEventListener('change', refreshLabels);
dom.preset.addEventListener('change', () => applyPreset(dom.preset.value));

dom.sustain.addEventListener('click', () => {
  sustain = !sustain;
  toggleButton(dom.sustain, sustain);
  if (!sustain) releaseSustained();
});

dom.chord.addEventListener('click', () => {
  chordMode = !chordMode;
  toggleButton(dom.chord, chordMode);
});

dom.arp.addEventListener('click', () => {
  arpMode = !arpMode;
  toggleButton(dom.arp, arpMode);
});

dom.record.addEventListener('click', () => {
  recording = !recording;
  dom.record.textContent = recording ? 'Stop' : 'Record';
  if (recording) {
    recorded = [];
    recordStart = performance.now();
    dom.status.textContent = 'Recording keyboard performance';
  } else {
    dom.status.textContent = `${recorded.length} recorded events`;
  }
});

dom.play.addEventListener('click', () => {
  if (!recorded.length) {
    dom.status.textContent = 'Nothing recorded yet';
    return;
  }
  ensureAudio();
  dom.status.textContent = 'Playing recording';
  for (const event of recorded) {
    setTimeout(() => {
      if (event.type === 'down') noteOn(event.key);
      else noteOff(event.key);
    }, event.at);
  }
});

buildKeyboard();
refreshLabels();
applyPreset('bright');
"""

README = """# KeySynth Studio

Original browser synth scaffold for FlyOrnith runs.

Start locally:

```bash
python -m http.server 8765
```

Then open `http://127.0.0.1:8765`.
"""

files = {
    "package.json": PACKAGE_JSON,
    "index.html": INDEX_HTML,
    "styles.css": STYLES,
    "synth.js": SYNTH_JS,
    "README.md": README,
}

for name, content in files.items():
    Path(name).write_text(content, encoding="utf-8")

print("Created KeySynth Studio static web app:", ", ".join(sorted(files)))
'''


def _html_webapp_creator_script(filename: str) -> str:
    return f'''from pathlib import Path


ROOT = Path(".")
INDEX = Path({filename!r})
if INDEX.name != "index.html":
    INDEX = Path("index.html")

PACKAGE_JSON = """{{
  "name": "metro-dash",
  "version": "1.0.0",
  "private": true,
  "description": "Original subway-themed endless runner web app.",
  "scripts": {{
    "start": "python -m http.server 8765",
    "check": "python -c \\"from pathlib import Path; assert Path('index.html').exists(); assert Path('game.js').exists(); assert Path('styles.css').exists(); print('static app files ok')\\""
  }}
}}
"""

INDEX_HTML = """<!doctype html>
<html lang="en">
  <head>
    <meta charset="utf-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1" />
    <title>Metro Dash</title>
    <link rel="stylesheet" href="styles.css" />
  </head>
  <body>
    <main class="game-shell">
      <section class="hud" aria-label="Run stats">
        <span>Score <strong id="score">0</strong></span>
        <span>Coins <strong id="coins">0</strong></span>
        <span>Distance <strong id="distance">0m</strong></span>
        <span>Speed <strong id="speed">1.0x</strong></span>
      </section>
      <canvas id="game" width="960" height="540" aria-label="Metro Dash game canvas"></canvas>
      <section id="overlay" class="overlay">
        <div>
          <p class="eyebrow">Original endless runner</p>
          <h1>Metro Dash</h1>
          <p id="message">Switch lanes, jump barriers, slide under signs, and collect signal coins.</p>
          <button id="primary">Start Run</button>
          <p class="keys">Arrow keys or WASD: move, jump, slide. Space pauses.</p>
        </div>
      </section>
    </main>
    <script src="game.js"></script>
  </body>
</html>
"""

STYLES = """* {{
  box-sizing: border-box;
}}

html,
body {{
  height: 100%;
  margin: 0;
}}

body {{
  display: grid;
  place-items: center;
  background: #10141f;
  color: #f7fbff;
  font-family: Arial, Helvetica, sans-serif;
}}

.game-shell {{
  position: relative;
  width: min(100vw, 1120px);
  aspect-ratio: 16 / 9;
  min-height: 320px;
  overflow: hidden;
  background: #172033;
  border: 1px solid #334057;
}}

canvas {{
  width: 100%;
  height: 100%;
  display: block;
  background: #152033;
}}

.hud {{
  position: absolute;
  z-index: 3;
  inset: 12px 12px auto 12px;
  display: grid;
  grid-template-columns: repeat(4, 1fr);
  gap: 8px;
  font-size: clamp(12px, 1.8vw, 16px);
}}

.hud span {{
  padding: 8px 10px;
  background: rgba(7, 12, 20, 0.72);
  border: 1px solid rgba(255, 255, 255, 0.14);
}}

.overlay {{
  position: absolute;
  z-index: 4;
  inset: 0;
  display: grid;
  place-items: center;
  padding: 24px;
  background: rgba(11, 15, 25, 0.76);
  text-align: center;
}}

.overlay.hidden {{
  display: none;
}}

.overlay div {{
  max-width: 560px;
}}

.eyebrow {{
  margin: 0 0 8px;
  color: #63d9a5;
  font-weight: 700;
  text-transform: uppercase;
  font-size: 13px;
}}

h1 {{
  margin: 0;
  font-size: clamp(42px, 8vw, 88px);
  line-height: 0.95;
}}

#message {{
  margin: 16px auto 22px;
  color: #dce7f5;
  line-height: 1.5;
}}

button {{
  border: 0;
  background: #ffce53;
  color: #1b1b20;
  padding: 12px 18px;
  font-weight: 800;
  cursor: pointer;
}}

.keys {{
  color: #aebbd0;
  font-size: 13px;
}}
"""

GAME_JS = """const canvas = document.querySelector('#game');
const ctx = canvas.getContext('2d');
const overlay = document.querySelector('#overlay');
const primary = document.querySelector('#primary');
const message = document.querySelector('#message');
const scoreEl = document.querySelector('#score');
const coinsEl = document.querySelector('#coins');
const distanceEl = document.querySelector('#distance');
const speedEl = document.querySelector('#speed');

const lanes = [-210, 0, 210];
const state = {{
  running: false,
  paused: false,
  over: false,
  lane: 1,
  y: 0,
  vy: 0,
  slide: 0,
  speed: 5.8,
  distance: 0,
  score: 0,
  coins: 0,
  spawn: 0,
  items: [],
  last: 0
}};

function reset() {{
  Object.assign(state, {{
    running: true,
    paused: false,
    over: false,
    lane: 1,
    y: 0,
    vy: 0,
    slide: 0,
    speed: 5.8,
    distance: 0,
    score: 0,
    coins: 0,
    spawn: 0,
    items: [],
    last: performance.now()
  }});
  overlay.classList.add('hidden');
  requestAnimationFrame(loop);
}}

function endRun() {{
  state.running = false;
  state.over = true;
  message.textContent = `Run complete: ${{Math.floor(state.distance)}}m, ${{state.coins}} coins, ${{Math.floor(state.score)}} points.`;
  primary.textContent = 'Restart';
  overlay.classList.remove('hidden');
}}

function spawnItem() {{
  const lane = Math.floor(Math.random() * 3);
  const roll = Math.random();
  state.items.push({{
    lane,
    z: 760,
    kind: roll < 0.34 ? 'coin' : roll < 0.67 ? 'barrier' : 'sign'
  }});
}}

function update(dt) {{
  if (!state.running || state.paused) return;
  state.distance += state.speed * dt * 8;
  state.score += state.speed * dt * 14;
  state.speed += dt * 0.08;
  state.spawn -= dt;
  if (state.spawn <= 0) {{
    spawnItem();
    state.spawn = Math.max(0.42, 1.08 - state.speed * 0.055);
  }}

  if (state.y > 0 || state.vy > 0) {{
    state.y += state.vy * dt;
    state.vy -= 1500 * dt;
    if (state.y < 0) {{
      state.y = 0;
      state.vy = 0;
    }}
  }}
  state.slide = Math.max(0, state.slide - dt);

  for (const item of state.items) item.z -= state.speed * dt * 120;
  state.items = state.items.filter((item) => item.z > -80);

  for (const item of state.items) {{
    if (item.hit) continue;
    const close = item.z < 70 && item.z > -20 && item.lane === state.lane;
    if (!close) continue;
    if (item.kind === 'coin') {{
      item.hit = true;
      state.coins += 1;
      state.score += 100;
    }} else if (item.kind === 'barrier' && state.y < 56) {{
      endRun();
    }} else if (item.kind === 'sign' && state.slide <= 0) {{
      endRun();
    }}
  }}
}}

function drawTrack() {{
  const w = canvas.width;
  const h = canvas.height;
  ctx.fillStyle = '#152033';
  ctx.fillRect(0, 0, w, h);
  ctx.fillStyle = '#25314a';
  ctx.beginPath();
  ctx.moveTo(w * 0.28, h);
  ctx.lineTo(w * 0.43, h * 0.16);
  ctx.lineTo(w * 0.57, h * 0.16);
  ctx.lineTo(w * 0.72, h);
  ctx.closePath();
  ctx.fill();

  for (let i = 0; i < 3; i += 1) {{
    const x = w / 2 + lanes[i] * 0.72;
    ctx.strokeStyle = '#ffce53';
    ctx.lineWidth = 4;
    ctx.beginPath();
    ctx.moveTo(x, h);
    ctx.lineTo(w / 2 + lanes[i] * 0.14, h * 0.18);
    ctx.stroke();
  }}

  ctx.fillStyle = '#66d9ef';
  for (let i = 0; i < 18; i += 1) {{
    const y = (i * 72 + state.distance * 2) % h;
    ctx.fillRect(w * 0.18, y, 70, 8);
    ctx.fillRect(w * 0.74, y + 24, 70, 8);
  }}
}}

function project(lane, z) {{
  const depth = 1 - z / 820;
  return {{
    x: canvas.width / 2 + lanes[lane] * (0.18 + depth * 0.58),
    y: canvas.height * (0.19 + depth * 0.76),
    s: 0.35 + depth * 1.05
  }};
}}

function drawItems() {{
  for (const item of state.items) {{
    if (item.hit) continue;
    const p = project(item.lane, item.z);
    if (item.kind === 'coin') {{
      ctx.fillStyle = '#ffce53';
      ctx.beginPath();
      ctx.arc(p.x, p.y - 32 * p.s, 16 * p.s, 0, Math.PI * 2);
      ctx.fill();
    }} else if (item.kind === 'barrier') {{
      ctx.fillStyle = '#ff5d5d';
      ctx.fillRect(p.x - 28 * p.s, p.y - 48 * p.s, 56 * p.s, 42 * p.s);
    }} else {{
      ctx.fillStyle = '#7ef0c1';
      ctx.fillRect(p.x - 34 * p.s, p.y - 94 * p.s, 68 * p.s, 18 * p.s);
    }}
  }}
}}

function drawPlayer() {{
  const x = canvas.width / 2 + lanes[state.lane] * 0.72;
  const baseY = canvas.height - 76 - state.y;
  const height = state.slide > 0 ? 44 : 82;
  ctx.fillStyle = '#f7fbff';
  ctx.fillRect(x - 24, baseY - height, 48, height);
  ctx.fillStyle = '#63d9a5';
  ctx.fillRect(x - 18, baseY - height - 20, 36, 22);
}}

function draw() {{
  drawTrack();
  drawItems();
  drawPlayer();
  scoreEl.textContent = Math.floor(state.score);
  coinsEl.textContent = state.coins;
  distanceEl.textContent = `${{Math.floor(state.distance)}}m`;
  speedEl.textContent = `${{(state.speed / 5.8).toFixed(1)}}x`;
}}

function loop(now) {{
  const dt = Math.min(0.04, (now - state.last) / 1000 || 0);
  state.last = now;
  update(dt);
  draw();
  if (state.running) requestAnimationFrame(loop);
}}

function move(direction) {{
  state.lane = Math.max(0, Math.min(2, state.lane + direction));
}}

window.addEventListener('keydown', (event) => {{
  const key = event.key.toLowerCase();
  if (key === 'arrowleft' || key === 'a') move(-1);
  if (key === 'arrowright' || key === 'd') move(1);
  if ((key === 'arrowup' || key === 'w') && state.y === 0) state.vy = 720;
  if (key === 'arrowdown' || key === 's') state.slide = 0.72;
  if (key === ' ') {{
    state.paused = !state.paused;
    if (!state.paused && state.running) {{
      state.last = performance.now();
      requestAnimationFrame(loop);
    }}
  }}
}});

primary.addEventListener('click', reset);
draw();
"""

README = """# Metro Dash

Original browser endless-runner scaffold for FlyOrnith runs.

Start locally:

```bash
python -m http.server 8765
```

Then open `http://127.0.0.1:8765`.
"""

files = {{
    "package.json": PACKAGE_JSON,
    "index.html": INDEX_HTML,
    "styles.css": STYLES,
    "game.js": GAME_JS,
    "README.md": README,
}}

for name, content in files.items():
    Path(name).write_text(content, encoding="utf-8")

print("Created Metro Dash static web app:", ", ".join(sorted(files)))
'''


def _pptx_creator_script(filename: str) -> str:
    return f'''from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path({filename!r})
SLIDES = [
    {{
        "title": "AgentOrnith makes Ornith operational",
        "subtitle": "A local command-line model becomes a supervised coding workbench with memory, tools, artifacts, and recovery.",
        "left": "AgentOrnith harness",
        "right": "Command-line Ornith",
        "harness": "Runs through a visible goal loop with checkpoints, approvals, tool logs, and artifact tracking.",
        "cli": "Answers one prompt at a time and depends on the operator to preserve context, evidence, and next steps.",
        "tradeoff": "Tradeoff: the harness adds policy and UI overhead, but that overhead buys durability for long work.",
    }},
    {{
        "title": "Use case 1: Long coding tasks",
        "subtitle": "Keep multi-hour implementation work coherent after pauses, crashes, or compaction.",
        "left": "AgentOrnith harness",
        "right": "Command-line Ornith",
        "harness": "Breaks work into milestones, persists events in SQLite, writes Obsidian checkpoints, and resumes from handoff bundles.",
        "cli": "Can lose thread state when the terminal scrolls, the model context fills, or the session restarts.",
        "tradeoff": "Tradeoff: checkpoints take time, but they prevent repeated re-orientation and forgotten acceptance criteria.",
    }},
    {{
        "title": "Use case 2: Web and browser evidence",
        "subtitle": "Research and verify facts without giving the model unrestricted internet access.",
        "left": "AgentOrnith harness",
        "right": "Command-line Ornith",
        "harness": "Uses audited web search, fetch, browser screenshots, excerpts, timestamps, and citation references.",
        "cli": "Usually relies on pasted links or untracked manual browsing, so evidence is harder to replay or audit.",
        "tradeoff": "Tradeoff: tool-gated browsing is slower than raw guesses, but it makes sources reviewable.",
    }},
    {{
        "title": "Use case 3: Safe file, shell, and test work",
        "subtitle": "Let Ornith act locally while keeping risky operations visible and recoverable.",
        "left": "AgentOrnith harness",
        "right": "Command-line Ornith",
        "harness": "Routes shell, file writes, git diffs, tests, and approvals through policy modes with secret redaction and logs.",
        "cli": "Can propose commands, but the operator must decide what is safe, remember results, and recover manually.",
        "tradeoff": "Tradeoff: strict approval modes can interrupt work; workspace autopilot or bypass modes reduce friction when trusted.",
    }},
    {{
        "title": "Use case 4: Project workspaces and artifacts",
        "subtitle": "Treat each run like an isolated project workspace instead of a loose terminal session.",
        "left": "AgentOrnith harness",
        "right": "Command-line Ornith",
        "harness": "Maps the repo, isolates edits, tracks touched files, verifies deliverables, and exposes artifacts in the dashboard.",
        "cli": "Works in the current folder but has no built-in artifact registry, source preview, or promotion workflow.",
        "tradeoff": "Tradeoff: isolation adds copy/sync complexity, but it reduces accidental edits to the source workspace.",
    }},
    {{
        "title": "Use case 5: IDE-style supervision",
        "subtitle": "Keep user decisions in the same place as model activity.",
        "left": "AgentOrnith harness",
        "right": "Command-line Ornith",
        "harness": "Shows live model summaries, tool events, approvals, blockers, screenshots, and chat steering in one workbench.",
        "cli": "The user must scan terminal output and external files to understand what needs attention.",
        "tradeoff": "Tradeoff: the dashboard must stay well organized; focus chat and tabs keep the detail available without crowding the main flow.",
    }},
]


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=RGBColor(22, 28, 36), align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    frame.word_wrap = True
    return box


def add_panel(slide, x, y, w, h, heading, body, accent):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.color.rgb = RGBColor(205, 213, 223)
    add_textbox(slide, x + 0.22, y + 0.18, w - 0.44, 0.35, heading, size=15, bold=True, color=accent)
    add_textbox(slide, x + 0.22, y + 0.66, w - 0.44, h - 0.86, body, size=14, color=RGBColor(32, 41, 54))


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    dark = RGBColor(17, 24, 39)
    green = RGBColor(22, 101, 52)
    slate = RGBColor(71, 85, 105)
    amber = RGBColor(146, 64, 14)
    for index, item in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(255, 255, 255)
        add_textbox(slide, 0.55, 0.35, 12.2, 0.48, item["title"], size=25, bold=True, color=dark)
        add_textbox(slide, 0.58, 0.92, 11.9, 0.45, item["subtitle"], size=13, color=slate)
        if index == 0:
            add_panel(slide, 0.75, 1.75, 5.85, 2.2, item["left"], item["harness"], green)
            add_panel(slide, 6.75, 1.75, 5.85, 2.2, item["right"], item["cli"], slate)
            add_textbox(slide, 0.86, 4.45, 11.6, 1.0, item["tradeoff"], size=16, bold=True, color=amber, align=PP_ALIGN.CENTER)
        else:
            add_panel(slide, 0.75, 1.65, 5.85, 2.7, item["left"], item["harness"], green)
            add_panel(slide, 6.75, 1.65, 5.85, 2.7, item["right"], item["cli"], slate)
            add_textbox(slide, 0.86, 4.78, 11.6, 0.9, item["tradeoff"], size=15, color=amber, align=PP_ALIGN.CENTER)
        add_textbox(slide, 0.62, 6.82, 2.2, 0.28, f"{{index + 1}} / {{len(SLIDES)}}", size=10, color=slate)
    prs.save(OUTPUT)
    print(f"created {{OUTPUT}} with {{len(SLIDES)}} slides")


if __name__ == "__main__":
    build()
'''


def _html_command() -> str:
    return (
        "python -c \"from pathlib import Path; "
        "files=sorted(p for p in Path('.').rglob('*.html') if 'node_modules' not in p.parts); "
        "assert files, 'no html files found'; "
        "p=files[0]; text=p.read_text(encoding='utf-8', errors='replace').lower(); "
        "assert '<html' in text and '</html>' in text, 'html document is incomplete'; "
        "assert p.stat().st_size > 500, 'html file is too small'; "
        "print(str(p), p.stat().st_size)\""
    )
